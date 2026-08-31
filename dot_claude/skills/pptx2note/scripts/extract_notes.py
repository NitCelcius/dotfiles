"""Extract speaker notes from every .pptx in a directory into individual _notes.txt files.

Usage:
    uv run --with python-pptx python extract_notes.py <directory>

Writes one <stem>_notes.txt next to each <stem>.pptx, containing:
    --- Slide 1 ---
    <notes text, or "(no notes)">

    --- Slide 2 ---
    ...

Files that fail to open (encrypted/password-protected, corrupt, not actually
a .pptx) are reported and skipped rather than aborting the whole batch.
"""
import sys
from pathlib import Path
from pptx import Presentation


def extract_all(src_dir: Path) -> None:
    pptx_files = sorted(p for p in src_dir.glob("*.pptx") if not p.name.startswith("~$"))

    if not pptx_files:
        print("No .pptx files found.")
        return

    for pptx_path in pptx_files:
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            print(f"FAILED to open {pptx_path.name}: {e}")
            continue

        lines = []
        for i, slide in enumerate(prs.slides, start=1):
            notes_text = ""
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                if tf is not None:
                    notes_text = tf.text.strip()
            lines.append(f"--- Slide {i} ---")
            lines.append(notes_text if notes_text else "(no notes)")
            lines.append("")

        out_path = pptx_path.with_name(pptx_path.stem + "_notes.txt")
        out_path.write_text("\n".join(lines), encoding="utf-8")
        print(f"Wrote {out_path.name} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: extract_notes.py <directory>")
        sys.exit(1)
    extract_all(Path(sys.argv[1]))
